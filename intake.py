"""Stage 1 — intake. A user request becomes a plan.md the Stage-2 designer
consumes. One unified path handles all three input situations: a vague
request, a request plus source documents, or an already-rough plan.

Source documents (PDF / DOCX / PPTX) are parsed with docling — IBM Research's
layout-aware document converter — which yields clean Markdown.
"""
from __future__ import annotations

import logging
import re
from pathlib import Path

import config
import llm
from harness_prompts import (
    CRAFTER_SYSTEM_PROMPT,
    TRANSCRIBE_SYSTEM_PROMPT,
    build_crafter_user_message,
    build_transcribe_user_message,
)

log = logging.getLogger("intake")

# Curated real training plans (reference_plans/) shown to the crafter as
# format + density exemplars. Chosen to span the range and, deliberately, to
# cover the hard high-value cases: multi-visual slides, matrix tables, and
# layered architecture diagrams.
EXEMPLAR_NAMES = [
    "multiviz_readout.md",        # multi-visual slides; bar / line / donut / funnel
    "competitive_overview.md",    # matrix table, 2x2 quadrant, brief-style discipline
    "architecture_reference.md",  # layered architecture, components, request flow
    "kafka_fundamentals.md",      # code blocks, stacked layers, definition
    "stripe_board_update.md",     # terse business deck -- stats, comparison
]

# Paired (source.txt, plan.md) exemplars for the TRANSCRIBE path -- shown to
# the crafter when the user uploads a finished slide deck and wants a 1:1
# plan that regenerates it. Each entry is a subdir under
# reference_plans/transcribe/ containing source.txt + plan.md.
TRANSCRIBE_PAIRS = [
    "01_credit_exception",   # multi-region one-slide (header / 3 circles / 2-col body)
    "02_hero_stat",          # simple one-slide (single hero number)
    "03_quarterly_readout",  # multi-slide deck, one ## per source slide
]

# Source kinds that imply "this is already a deck, transcribe it" rather than
# "this is raw research material, organize it." Add .key here once we test it.
_DECK_SUFFIXES = {".pptx", ".key"}

_FENCE = re.compile(r"^\s*```[a-zA-Z]*\s*\n(.*?)\n?\s*```\s*$", re.DOTALL)

# The crafter is told to use plain ASCII, but the model still emits the odd
# typographic dash / curly quote. Those corrupt downstream (the `?`/`�`
# bug), so normalize deterministically rather than trusting the instruction.
_ASCII_MAP = str.maketrans({
    "–": "-", "—": "-", "‒": "-", "‐": "-", "‑": "-",
    "―": "-", "−": "-",
    "“": '"', "”": '"', "„": '"', "‘": "'",
    "’": "'", "‚": "'",
    "…": "...", "•": "-", " ": " ", "×": "x",
})


_ASCII_MAP.update(str.maketrans({"≈": "~", "≠": "!=", "±": "+/-"}))


def _strip_fence(text: str) -> str:
    m = _FENCE.match(text.strip())
    return m.group(1).strip() if m else text.strip()


def _normalize_ascii(text: str) -> str:
    """Fold typographic punctuation the crafter occasionally emits down to
    plain ASCII (en/em dashes, curly quotes, ellipsis) — those do not survive
    the designer -> coder -> render pipeline cleanly."""
    out = text.translate(_ASCII_MAP)
    out = out.replace("≥", ">=").replace("≤", "<=")
    out = out.replace("→", "->")
    return out


def _load_exemplars() -> list[str]:
    out: list[str] = []
    for name in EXEMPLAR_NAMES:
        p = config.REFERENCE_PLANS / name
        if p.is_file():
            out.append(p.read_text())
    return out


def _load_transcribe_pairs() -> list[tuple[str, str]]:
    """Load (source.txt, plan.md) exemplar pairs for the TRANSCRIBE path. A
    missing file in a pair drops that pair silently -- exemplar curation
    breakage shouldn't abort intake."""
    out: list[tuple[str, str]] = []
    base = config.REFERENCE_PLANS / "transcribe"
    for d in TRANSCRIBE_PAIRS:
        sp = base / d / "source.txt"
        pp = base / d / "plan.md"
        if sp.is_file() and pp.is_file():
            out.append((sp.read_text(), pp.read_text()))
        else:
            log.warning("transcribe exemplar missing: %s", d)
    return out


def _is_deck_source(paths: list[Path]) -> bool:
    """Heuristic: any uploaded .pptx (or .key) triggers transcribe mode. PDFs
    are NOT treated as decks by default -- a research paper looks the same as
    a slide-PDF to the suffix check, so PDFs stay on the research path. A
    user who wants transcribe behaviour on a slide-PDF can convert to PPTX
    or, later, we add a slide-density heuristic."""
    return any(p.suffix.lower() in _DECK_SUFFIXES for p in paths)


def _pptx_to_text(path: Path) -> str:
    """Walk a PPTX slide-by-slide and emit text with explicit
    `--- slide N ---` boundaries. Docling flattens decks into a single
    Markdown stream with no slide boundaries, which makes the transcribe
    crafter unable to count source slides; this path preserves them.

    Reading order within a slide is left-to-right, top-to-bottom by shape
    position -- close enough for the transcribe crafter's purposes, and
    matches what the model expects from exemplar 3."""
    from pptx import Presentation
    prs = Presentation(str(path))
    parts: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        parts.append(f"--- slide {i} ---")
        shapes = sorted(
            (s for s in slide.shapes if getattr(s, "top", None) is not None),
            key=lambda s: (int((s.top or 0) // 100000),  # row band
                           int((s.left or 0) // 100000)),  # then column
        )
        for sh in shapes:
            if sh.has_text_frame:
                for para in sh.text_frame.paragraphs:
                    line = "".join(r.text for r in para.runs).strip()
                    if line:
                        parts.append(line)
            elif getattr(sh, "shape_type", None) == 19:  # TABLE
                for row in sh.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    parts.append(" | ".join(cells))
        parts.append("")  # blank line between slides
    return "\n".join(parts).strip()


def parse_document(path: Path) -> tuple[str, str]:
    """Any supported document -> (filename, markdown text).

    PPTX uses a slide-aware walker that emits `--- slide N ---` boundaries
    so the transcribe crafter can count source slides. PDF / DOCX go through
    docling (layout-aware, outputs Markdown with reading order and tables
    preserved). Plain .md / .txt are read directly.
    """
    path = Path(path)
    suffix = path.suffix.lower()
    if suffix in (".md", ".txt"):
        return path.name, path.read_text(errors="replace")
    if suffix == ".pptx":
        return path.name, _pptx_to_text(path)
    from docling.document_converter import DocumentConverter
    result = DocumentConverter().convert(str(path))
    return path.name, result.document.export_to_markdown()


def craft_plan(request: str, source_paths: list[Path] | None = None) -> str:
    """Turn a request (plus any uploaded source documents) into a plan.md."""
    source_texts: list[tuple[str, str]] = []
    for sp in source_paths or []:
        try:
            source_texts.append(parse_document(sp))
            log.info("parsed source: %s", Path(sp).name)
        except Exception as e:  # noqa: BLE001 — a bad upload shouldn't abort
            log.warning("failed to parse %s: %s", sp, e)

    transcribe = _is_deck_source(source_paths or [])
    if transcribe:
        log.info("transcribe mode: uploaded source is a finished deck")
        messages = [
            {"role": "system", "content": TRANSCRIBE_SYSTEM_PROMPT},
            {"role": "user", "content": build_transcribe_user_message(
                request, source_texts, _load_transcribe_pairs())},
        ]
    else:
        messages = [
            {"role": "system", "content": CRAFTER_SYSTEM_PROMPT},
            {"role": "user", "content": build_crafter_user_message(
                request, source_texts, _load_exemplars())},
        ]
    content, _ = llm.chat(config.ROSTER["crafter"], messages)
    plan = _normalize_ascii(_strip_fence(content))
    log.info("crafted plan (%s): %d chars, %d slide section(s)",
             "transcribe" if transcribe else "research",
             len(plan), plan.count("\n## "))
    return plan
